from fastapi.responses import FileResponse, HTMLResponse
from python_pptx_text_replacer import TextReplacer
from pptx.chart.data import CategoryChartData
from pptx import Presentation
from fastapi import FastAPI




app = FastAPI()


# Fungsi untuk mengganti data dalam chart pada slide tertentu
def replace_chart_with_data(slide, chart_index, chart_data):
    chart_count = 0
    for shape in slide.shapes:
        if shape.has_chart:
            chart_count += 1
            if chart_count == chart_index + 1:  
                chart = shape.chart
                chart.replace_data(chart_data)
                print(f"Chart with index {chart_index} found and replaced successfully on the specified slide.")
                return
    print(f"Chart with index {chart_index} not found on the specified slide.")



# Endpoint untuk generate
@app.post("/generate")
async def generate(Any: dict):
    
    # membuka pptx untuk di lakukan replace Chart
    prs = Presentation('File Lama/Pitching_Report_DUPK_BI_Pengaduan_Sistem_Pembayaran_1_7_April_2022.pptx')

    # data top_10_online_media
    data = Any['result']['top_10_online_media']
    sorted_data = sorted(data.items(), key=lambda x: x[1], reverse=False)
    categories = [item[0] for item in sorted_data]
    values = [item[1] for item in sorted_data]

    
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)

    
    slide_index = 2
    chart_index_to_replace = 0

    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break

    
    # data top_10_printed_media
    data = Any['result']['top_10_printed_media']
    sorted_data = sorted(data.items(), key=lambda x: x[1], reverse=False)
    categories = [item[0] for item in sorted_data]
    values = [item[1] for item in sorted_data]

    
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)

    
    slide_index = 2
    chart_index_to_replace = 1

    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break
    


    # data per_day_detail
    data = Any['result']['per_day_detail']
    dates = []
    online_values = []
    printed_values = []

    for date, details in data.items():
        dates.append(date)
        online_values.append(details['online'])
        printed_values.append(details['printed'])

    
    chart_data = CategoryChartData()
    chart_data.categories = dates
    chart_data.add_series('', online_values)
    chart_data.add_series('', printed_values)

    slide_index = 2
    chart_index_to_replace = 2

    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break



    # data top_10_online_influencer
    data = Any['result']['top_10_online_influencer']
    sorted_data = sorted(data.items(), key=lambda x: x[1], reverse=False)
    categories = [item[0] for item in sorted_data]
    values = [item[1] for item in sorted_data]

    
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    
    slide_index = 4
    chart_index_to_replace = 0

    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break



    # Data all_days_detail
    data = Any['result']['all_days_detail']
    categories = [item['text'] for item in data]
    totals = [item['percentage'] for item in data]

    
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', totals)

    slide_index = 3
    chart_index_to_replace = 0

    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break



    # Data isu dan Sentiment
    data = Any['result']['sentiment']
    categories = ["Negatif" ,"Netral" ,"Positif"]
    totals = [data['negative']['percentage'],data['neutral']['percentage'],data['positive']['percentage']]


    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', totals)

    slide_index = 3
    chart_index_to_replace = 1

    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break

    

    # top_10_printed_influencer
    data = Any['result']['top_10_printed_influencer']
    sorted_data = sorted(data.items(), key=lambda x: x[1], reverse=False)
    categories = [item[0] for item in sorted_data]
    values = [item[1] for item in sorted_data]

    
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)

    slide_index = 4
    chart_index_to_replace = 1

    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, chart_data)
            break
    

    # menyimpan data dahulu setelah selesai mereplace chart, dan selanjutnya data yang sudah di replace akan dibaca kembali untuk di replace teksnya
    save_file = ('Template/Pitching_Report_DUPK_BI_Pengaduan_Sistem_Pembayaran_1_7_April_2022.pptx')
    prs.save(save_file)

    data = Any['result']
    online_news = data['total_online_news']
    online_media = data['total_online_media']
    printed_news = data['total_printed_news']
    printed_media = data['total_printed_media']
    

    # Ambil data topik dari tanggal tertentu
    # data untuk mengganti di pptx slide ke - 3 dan diambil satu satu melewati index
    data_perday = data['per_day_detail']
    data_teks1 = [topic['text'] for topic in data_perday['2022-04-01']['topics']]
    data_teks2 = [topic['text'] for topic in data_perday['2022-04-07']['topics']]
    # tanggal 2022-04-01
    slide3_1 = data_teks1[0]
    slide3_2 = data_teks1[1]
    slide3_3 = data_teks1[2]
    slide3_4 = data_teks1[3]
    # 2022-04-07
    slide3_5 = data_teks2[0]
    slide3_6 = data_teks2[3]
    slide3_7 = data_teks2[2]
    
    # data untuk mengganti di pptx slide ke - 4 dan diambil satu satu melewati index
    data_isu_sentiment = data['sentiment']
    data_teks3 = data_isu_sentiment['negative']['topic_examples']
    data_teks4 = data_isu_sentiment['positive']['topic_examples']

    # untuk data isu dan sentiment negatif
    slide4_1 = data_teks3[0]
    slide4_2 = data_teks3[1]
    slide4_3 = data_teks3[2]
    slide4_4 = data_teks3[3]

    # untuk data isu dan sentiment positif
    slide4_5 = data_teks4[0]
    slide4_6 = data_teks4[1]
    slide4_7 = data_teks4[2]
    slide4_8 = data_teks4[3]

    # Membuat variabel TextReplacer untuk mengganti teks dalam presentasi
    replacer = TextReplacer(save_file, slides='', tables=False, charts=False, textframes=True)

    # Mengganti teks dalam presentasi
    replacer.replace_text([
        ('2.251', str(online_news)),
        ('761', str(online_media)),
        ('135', str(printed_news)),
        ('60', str(printed_media)),
        ('1 – 7 April 2022' , str(data['earliest_date'] + " Sampai " + str(data['latest_date']))),

        # untuk tanggal 2022-04-01
        ('Gangguan layanan mobile banking BCA', slide3_1),
        ('Perluasan implementasi QRIS ', slide3_2),
        ('BPJPH kembangkan Sistem Informasi Halal (Sihalal) yang terintegrasi dengan penyedia uang elektronik', slide3_3),
        ('Promo belanja menggunakan kartu debit dan kredit serta dompet digital', slide3_4),

        # untuk tanggal 2022-04-07
        ('BI dorong penggunaan transaksi nontunai selama Ramadan dan Idulfitri', slide3_5),
        ('Fitur Digital dalam Bulan Ramadhan.', slide3_6),
        ('Pemerintah berencana memajaki fintech dan dompet digital', slide3_7),

        # untuk data isu dan sentiment negatif
        ('Keluhan masyarakat perihal gangguan pada layanan mobile banking BCA. [link]', slide4_1),
        ('Terungkapnya modus skimming melalui modus pengganjal ATM di Cilacap. [link]', slide4_2),
        ('Terungkapnya dugaan kasus skimming nasabah BNI di Samarinda. [link]', slide4_3),
        ('Pencatutan identitas sebabkan kerugian berupa kesulitan pengajuan kartu kredit. [link]', slide4_4),
        ('Keluhan soal saldo yang tidak kunjung masuk meskipun proses scan QRIS sudah berhasil. [link]', ''),
        ('Ketimpangan penyaluran pinjaman online antara Pulau Jawa dan wilayah lainnya. [link]', ''),


        # untuk data isu dan sentiment positif
        ('BI mendorong perluasan transaksi nontunai di masyarakat.', slide4_5),
        ('Kontribusi perbankan, penyedia dompet digital, dan pemerintah mendorong transaksi nontunai.', slide4_6),
        ('Perbankan pastikan keamanan jaringan untuk transaksi di ATM selama Ramadan dan Idulfitri.', slide4_7),
        ('Pemerintah berencana memudahkan transaksi Pemda melalui Kartu Kredit Pemerintah Daerah (KKPD).', slide4_8)
    ])

    # Menyimpan presentasi yang telah diubah
    file_output = "File Baru/Pitching_Report_DUPK_BI_Pengaduan_Sistem_Pembayaran_1_7_April_2022.pptx"
    replacer.write_presentation_to_file(file_output)

    # Mengembalikan file yang dapat diunduh
    return FileResponse(file_output, filename="Pitching_Report_DUPK_BI_Pengaduan_Sistem_Pembayaran_1_7_April_2022.pptx")
